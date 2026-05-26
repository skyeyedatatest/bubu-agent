from typing import Dict, Any, List, Optional, Tuple
from pathlib import Path
import re
import urllib.request
import tempfile
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import PP_PLACEHOLDER

class SkillTools:
    """
    PPT 生成技能实现类 (Enhanced v2)
    """
    
    def __init__(self, workspace_root: str):
        self.workspace_root = Path(workspace_root)
        self._ensure_workspace()

    def _ensure_workspace(self):
        if not self.workspace_root.exists():
            self.workspace_root.mkdir(parents=True, exist_ok=True)

    def execute(self, params: Dict[str, Any]) -> str:
        """
        执行入口
        """
        markdown_content = params.get("markdown_content", "")
        filename = params.get("filename", "output.pptx")
        
        if not markdown_content:
            return "❌ 错误: 请提供 Markdown 内容大纲。"
            
        if not filename.endswith(".pptx"):
            filename += ".pptx"
            
        try:
            output_path = self.workspace_root / filename
            self._generate_pptx(markdown_content, output_path)
            
            return f"✅ PPT 已成功生成！\n\n文件位置: {output_path}\n您可以在 Data 目录中找到它。"
            
        except Exception as e:
            return f"❌ 生成失败: {str(e)}"

    def _generate_pptx(self, md_content: str, output_path: Path):
        """
        核心生成逻辑：Markdown -> PPTX (Smart Layout Engine)
        """
        # Load template if exists
        template_path = Path(__file__).parent / "template.pptx"
        if template_path.exists():
            prs = Presentation(str(template_path))
        else:
            prs = Presentation()
        
        # 1. Parse Markdown into Slide Blocks
        slides_data = self._parse_markdown_to_blocks(md_content)
        
        # 2. Process each slide
        for i, slide_data in enumerate(slides_data):
            self._create_slide(prs, slide_data, is_first=(i==0))

        prs.save(str(output_path))

    def _parse_markdown_to_blocks(self, md_content: str) -> List[Dict]:
        """
        解析 Markdown 为结构化数据块
        """
        blocks = []
        lines = md_content.split('\n')
        current_block = {"title": "", "content": [], "images": [], "layout": None}
        
        img_pattern = re.compile(r'!\[(.*?)\]\((.*?)\)')
        layout_pattern = re.compile(r'<!--\s*layout:\s*(.*?)\s*-->', re.IGNORECASE)
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            # Check for Layout Comment
            layout_match = layout_pattern.match(line)
            if layout_match:
                current_block["layout"] = layout_match.group(1)
                continue

            # Check for Title (Start of new slide)
            if line.startswith('# '):
                # Save previous block if valid
                if current_block["title"] or current_block["content"] or current_block["images"]:
                    blocks.append(current_block)
                
                # Start new block
                current_block = {
                    "title": line[2:].strip(),
                    "content": [],
                    "images": [],
                    "layout": None
                }
                continue
            
            # Check for Images
            img_match = img_pattern.match(line)
            if img_match:
                current_block["images"].append({
                    "alt": img_match.group(1),
                    "url": img_match.group(2)
                })
                continue
                
            # Regular Content
            current_block["content"].append(line)
            
        # Append last block
        if current_block["title"] or current_block["content"] or current_block["images"]:
            blocks.append(current_block)
            
        return blocks

    def _create_slide(self, prs: Presentation, slide_data: Dict, is_first: bool):
        """
        根据数据创建单个幻灯片
        """
        title = slide_data["title"]
        content_lines = slide_data["content"]
        images = slide_data["images"]
        layout_pref = slide_data["layout"]
        
        # --- Layout Selection Strategy ---
        layout_index = 1 # Default: Title and Content
        
        # Strategy 1: Explicit Layout
        if layout_pref:
            # Try to find layout by name
            for i, layout in enumerate(prs.slide_layouts):
                if layout.name.lower() == layout_pref.lower():
                    layout_index = i
                    break
        
        # Strategy 2: First Slide -> Title Slide (0)
        elif is_first:
            layout_index = 0
            
        # Strategy 3: Has Images -> Picture with Caption (8) or Two Content (3)
        elif images:
            # Prefer 'Picture with Caption' (8) if available
            if len(prs.slide_layouts) > 8: 
                layout_index = 8
            # Fallback to 'Two Content' (3)
            elif len(prs.slide_layouts) > 3:
                layout_index = 3
        
        # Bounds Check
        if layout_index >= len(prs.slide_layouts):
            layout_index = 1 if len(prs.slide_layouts) > 1 else 0

        slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
        
        # --- Content Population ---
        
        # 1. Set Title
        if slide.shapes.title:
            slide.shapes.title.text = title
            
        # 2. Identify Placeholders
        body_ph = None
        picture_ph = None
        
        # Find Body and Picture placeholders
        for shape in slide.placeholders:
            if shape == slide.shapes.title:
                continue
                
            try:
                ph_type = shape.placeholder_format.type
                if ph_type == PP_PLACEHOLDER.BODY or ph_type == PP_PLACEHOLDER.OBJECT:
                    if not body_ph: body_ph = shape
                    # If we have multiple bodies (Two Content), maybe use second for image?
                    elif not picture_ph and images: picture_ph = shape
                elif ph_type == PP_PLACEHOLDER.PICTURE:
                    picture_ph = shape
            except:
                pass
                
        # Fallback: If no explicit picture placeholder but we have an image, 
        # try to use the second body placeholder (for Two Content layout)
        if images and not picture_ph and layout_index == 3: # Two Content
             # Usually index 1 is body 1, index 2 is body 2
             # But let's trust the loop above finding the second body
             pass

        # 3. Fill Body Text
        target_body = body_ph if body_ph else None
        # If 'Picture with Caption', the body placeholder is usually the text one.
        
        if target_body and target_body.has_text_frame:
            tf = target_body.text_frame
            tf.clear() # Clear default prompt text
            
            for line in content_lines:
                if line.startswith('## '):
                    p = tf.add_paragraph()
                    p.text = line[3:].strip()
                    p.font.bold = True
                    p.font.size = Pt(20)
                elif line.startswith('- ') or line.startswith('* '):
                    p = tf.add_paragraph()
                    p.text = line[2:].strip()
                    p.level = 0
                elif line.startswith('  - ') or line.startswith('  * '):
                    p = tf.add_paragraph()
                    p.text = line[4:].strip()
                    p.level = 1
                else:
                    p = tf.add_paragraph()
                    p.text = line

        # 4. Fill Images
        if images:
            img_url = images[0]["url"] # Only handle first image for now
            
            if picture_ph:
                # Use placeholder
                self._add_image_to_placeholder(slide, picture_ph, img_url)
            else:
                # Manual placement (Fallback)
                self._add_image_manual(slide, img_url)

    def _add_image_to_placeholder(self, slide, placeholder, url):
        """Insert image into a specific placeholder"""
        image_path = self._download_image(url)
        if image_path:
            try:
                # Try standard insert
                if hasattr(placeholder, 'insert_picture'):
                    placeholder.insert_picture(image_path)
                else:
                    raise AttributeError("No insert_picture method")
                    
                self._cleanup_temp_image(image_path, url)
            except (AttributeError, Exception) as e:
                # Fallback: Manually place image using placeholder dimensions
                try:
                    left = placeholder.left
                    top = placeholder.top
                    width = placeholder.width
                    height = placeholder.height
                    
                    slide.shapes.add_picture(image_path, left, top, width=width, height=height)
                    self._cleanup_temp_image(image_path, url)
                except Exception as inner_e:
                    print(f"Failed to insert picture (fallback): {inner_e}")

    def _add_image_manual(self, slide, url):
        """Manually place image on the right"""
        image_path = self._download_image(url)
        if image_path:
            try:
                left = Inches(6.5)
                top = Inches(1.5)
                width = Inches(3.0)
                slide.shapes.add_picture(image_path, left, top, width=width)
                self._cleanup_temp_image(image_path, url)
            except Exception:
                pass

    def _download_image(self, url):
        if not url.startswith('http'):
            return url # Local path
            
        try:
            req = urllib.request.Request(url, headers={'User-Agent': 'Mozilla/5.0'})
            with urllib.request.urlopen(req, timeout=10) as response:
                data = response.read()
                ext = ".jpg"
                if ".png" in url.lower(): ext = ".png"
                elif ".gif" in url.lower(): ext = ".gif"
                
                fd, temp_path = tempfile.mkstemp(suffix=ext)
                os.close(fd)
                with open(temp_path, 'wb') as f:
                    f.write(data)
                return temp_path
        except Exception:
            return None

    def _cleanup_temp_image(self, path, original_url):
        if path != original_url and os.path.exists(path):
            try:
                os.remove(path)
            except:
                pass

